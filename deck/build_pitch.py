"""
KIRIM — the short pitch. Six slides, built to be presented in three minutes.

    python deck/build_pitch.py

Retia Labs design system. Icons are drawn from primitives rather than pulled
from a set, so they sit in the same geometric language as the rest of the deck —
hairlines, sharp corners, one accent.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
import os

# ---------------------------------------------------------------- Retia tokens
PAPER_0 = RGBColor(0xFA, 0xF6, 0xEE)
PAPER_1 = RGBColor(0xF1, 0xEC, 0xE1)
PAPER_2 = RGBColor(0xE8, 0xE2, 0xD4)
PAPER_3 = RGBColor(0xDE, 0xD7, 0xC6)
INK_0 = RGBColor(0x17, 0x13, 0x0D)
INK_2 = RGBColor(0x51, 0x4A, 0x3C)
INK_3 = RGBColor(0x6A, 0x63, 0x56)
INK_4 = RGBColor(0x8C, 0x84, 0x67)
PANEL_0 = RGBColor(0x0F, 0x0D, 0x09)
PANEL_1 = RGBColor(0x1A, 0x17, 0x12)
PANEL_2 = RGBColor(0x26, 0x21, 0x19)
PANEL_INK = RGBColor(0xD9, 0xD3, 0xC5)
PANEL_MUTED = RGBColor(0x9B, 0x94, 0x83)
PANEL_FAINT = RGBColor(0x7D, 0x76, 0x6A)
TRUSTED = RGBColor(0x2C, 0x5A, 0x4E)
CONTESTED = RGBColor(0xB0, 0x7A, 0x22)
QUARANTINED = RGBColor(0xA2, 0x3B, 0x2C)
RESTRICTED = RGBColor(0x7A, 0x2E, 0x2E)
SAGE = RGBColor(0x5E, 0x9B, 0x79)
TRUSTED_TINT = RGBColor(0xDC, 0xE6, 0xE0)
CONTESTED_TINT = RGBColor(0xEF, 0xE4, 0xCB)
QUARANTINED_TINT = RGBColor(0xEF, 0xDA, 0xD2)

DISPLAY = "Georgia"
SANS = "Segoe UI"
MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.9)
CW = W - 2 * M
LOGO_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", ".github", "logos")


# ---------------------------------------------------------------- primitives
def letterspace(run, ems):
    run._r.get_or_add_rPr().set("spc", str(int(ems * 100)))


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def write(tf, text, *, font=SANS, size=15, color=INK_0, bold=False, italic=False,
          spacing=0.0, line=1.35, space_after=0, first=False, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.line_spacing = line
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing:
        letterspace(r, spacing)
    return p


def box(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    sh.shadow.inherit = False
    return sh


def rule(slide, x, y, w, color=PAPER_3, weight=Pt(0.9)):
    sh = box(slide, x, y, w, Emu(int(weight.emu / 2)), fill=color)
    return sh


def slide_base(prs, bg=PAPER_1):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = box(s, 0, 0, W, H, fill=bg)
    bg_shape._element.getparent().remove(bg_shape._element)
    s.shapes._spTree.insert(2, bg_shape._element)
    return s


def eyebrow(slide, text, y=Inches(0.6), color=INK_4, x=M):
    tf = textbox(slide, x, y, CW, Inches(0.25))
    write(tf, text.upper(), font=MONO, size=10.5, color=color, spacing=0.18, first=True)


def footer(slide, left, right, color=INK_4, line=PAPER_3):
    rule(slide, M, H - Inches(0.7), CW, color=line)
    tf = textbox(slide, M, H - Inches(0.55), CW / 2, Inches(0.3))
    write(tf, left, font=MONO, size=9.5, color=color, spacing=0.12, first=True)
    tf2 = textbox(slide, M + CW / 2, H - Inches(0.55), CW / 2, Inches(0.3), align=PP_ALIGN.RIGHT)
    write(tf2, right, font=MONO, size=9.5, color=color, spacing=0.12, first=True, align=PP_ALIGN.RIGHT)


def logos(slide, top, right=None, dark=False):
    right = right or (W - M)
    pic = slide.shapes.add_picture(os.path.join(LOGO_DIR, "ripple.png"), Emu(0), top, height=Inches(0.28))
    pic.left = Emu(int(right - pic.width))
    div = box(slide, Emu(int(pic.left - Inches(0.26))), top - Inches(0.02),
              Emu(9525), Inches(0.32), fill=PANEL_2 if dark else PAPER_3)
    pic2 = slide.shapes.add_picture(os.path.join(LOGO_DIR, "singhacks.png"), Emu(0),
                                    top - Inches(0.02), height=Inches(0.32))
    pic2.left = Emu(int(pic.left - Inches(0.52) - pic2.width))
    tf = textbox(slide, Emu(int(pic2.left - Inches(1.3))), top + Inches(0.02),
                 Inches(1.1), Inches(0.3), align=PP_ALIGN.RIGHT)
    write(tf, "BUILT FOR", font=MONO, size=9, color=PANEL_FAINT if dark else INK_4,
          spacing=0.16, first=True, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- icons
# Drawn from primitives so they share the deck's geometry: hairline strokes,
# sharp corners, one accent colour. No icon set, no emoji.
def _stroke(slide, x, y, w, h, colour, weight=Inches(0.035)):
    return box(slide, x, y, w, h, fill=colour)


def icon_building(slide, x, y, size, colour):
    """Tower under construction: a stack of floors and a crane arm."""
    u = size / 10
    box(slide, x + u * 2, y + u * 3, u * 5, u * 7, fill=None, line_color=colour, line_w=Pt(1.6))
    for i in range(3):
        _stroke(slide, x + u * 2, y + u * (4.6 + i * 1.6), u * 5, Emu(int(u * 0.09)), colour)
    # crane mast and jib
    _stroke(slide, x + u * 8, y + u * 0.6, Emu(int(u * 0.12)), u * 9.4, colour)
    _stroke(slide, x + u * 3.4, y + u * 0.6, u * 5.4, Emu(int(u * 0.12)), colour)
    _stroke(slide, x + u * 4.2, y + u * 0.6, Emu(int(u * 0.1)), u * 1.5, colour)


def icon_lock(slide, x, y, size, colour):
    """Escrow: a closed lock."""
    u = size / 10
    box(slide, x + u * 1.6, y + u * 4.4, u * 6.8, u * 5, fill=None, line_color=colour, line_w=Pt(1.6))
    shackle = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, x + u * 2.6, y + u * 1.2, u * 4.8, u * 5.2)
    shackle.fill.background()
    shackle.line.color.rgb = colour
    shackle.line.width = Pt(1.6)
    shackle.shadow.inherit = False
    _stroke(slide, x + u * 4.85, y + u * 6.1, Emu(int(u * 0.3)), u * 1.6, colour)


def icon_shield(slide, x, y, size, colour):
    """Trust: a shield with a check."""
    u = size / 10
    sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, x + u * 1.5, y + u * 1.2, u * 7, u * 7.6)
    sh.rotation = 90
    sh.fill.background()
    sh.line.color.rgb = colour
    sh.line.width = Pt(1.6)
    sh.shadow.inherit = False
    tick1 = _stroke(slide, x + u * 3.6, y + u * 5.4, u * 1.9, Emu(int(u * 0.32)), colour)
    tick1.rotation = 45
    tick2 = _stroke(slide, x + u * 4.6, y + u * 4.4, u * 3.2, Emu(int(u * 0.32)), colour)
    tick2.rotation = -45


def icon_camera(slide, x, y, size, colour):
    """Evidence: a photograph with a location pin."""
    u = size / 10
    box(slide, x + u * 1.2, y + u * 2.8, u * 7.6, u * 5.4, fill=None, line_color=colour, line_w=Pt(1.6))
    _stroke(slide, x + u * 3.4, y + u * 1.9, u * 3.2, Emu(int(u * 0.12)), colour)
    lens = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + u * 3.7, y + u * 4.1, u * 2.6, u * 2.6)
    lens.fill.background()
    lens.line.color.rgb = colour
    lens.line.width = Pt(1.4)
    lens.shadow.inherit = False


def icon_coins(slide, x, y, size, colour):
    """Payment: stacked value moving."""
    u = size / 10
    for i, wdt in enumerate([6.4, 5.2, 4.0]):
        e = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + u * (1.6 + (6.4 - wdt) / 2),
                                   y + u * (6.6 - i * 2.0), u * wdt, u * 1.5)
        e.fill.background()
        e.line.color.rgb = colour
        e.line.width = Pt(1.4)
        e.shadow.inherit = False


def icon_clock(slide, x, y, size, colour):
    """Time: what everyone is waiting for."""
    u = size / 10
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + u * 1.2, y + u * 1.2, u * 7.6, u * 7.6)
    c.fill.background()
    c.line.color.rgb = colour
    c.line.width = Pt(1.6)
    c.shadow.inherit = False
    _stroke(slide, x + u * 4.9, y + u * 2.6, Emu(int(u * 0.26)), u * 2.6, colour)
    _stroke(slide, x + u * 5.0, y + u * 4.9, u * 2.2, Emu(int(u * 0.26)), colour)


# ---------------------------------------------------------------- the pitch
def flow_step(s, x, y, w, h, who, what, *, tone="neutral", tag=None):
    """One box in a vertical flow. Tone carries the argument, not decoration."""
    fill, line, who_c, what_c = {
        "neutral": (PAPER_0, PAPER_3, INK_0, INK_3),
        "bad": (QUARANTINED_TINT, QUARANTINED, QUARANTINED, RESTRICTED),
        "good": (TRUSTED_TINT, TRUSTED, TRUSTED, INK_2),
        "same": (PAPER_1, PAPER_3, INK_4, INK_4),
    }[tone]
    box(s, x, y, w, h, fill=fill, line_color=line)
    tf = textbox(s, x + Inches(0.22), y + Inches(0.11), w - Inches(1.3), h - Inches(0.2))
    write(tf, who, font=SANS, size=12.5, color=who_c, bold=True, first=True, space_after=1)
    write(tf, what, font=SANS, size=11, color=what_c, line=1.25)
    if tag:
        tf2 = textbox(s, x + w - Inches(1.15), y + Inches(0.13), Inches(0.95), Inches(0.25),
                      align=PP_ALIGN.RIGHT)
        write(tf2, tag, font=MONO, size=8, color=INK_4, spacing=0.1, first=True,
              align=PP_ALIGN.RIGHT)


def arrow(s, cx, y, colour=PAPER_3):
    box(s, cx - Emu(int(Pt(0.5).emu)), y, Emu(int(Pt(1).emu)), Inches(0.15), fill=colour)


def build(hashes):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ============================================================ 1 — title
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "SingHacks 2026 · Ripple · AI-native business on XRPL", color=PANEL_FAINT)
    rule(s, M, Inches(1.05), CW, color=PANEL_2)

    tf = textbox(s, M, Inches(2.35), Inches(9.6), Inches(1.6))
    write(tf, "Kirim", font=DISPLAY, size=76, color=PANEL_INK, line=1.0, first=True)
    tf2 = textbox(s, M, Inches(3.5), Inches(9.6), Inches(0.6))
    write(tf2, "to send", font=DISPLAY, size=26, color=CONTESTED, italic=True, line=1.1, first=True)

    rule(s, M, Inches(4.35), Inches(4.0), color=PANEL_2)
    tf = textbox(s, M, Inches(4.6), Inches(9.0), Inches(0.9))
    write(tf, "AI-verified milestone payments on the XRP Ledger.",
          font=SANS, size=19, color=PANEL_MUTED, line=1.4, first=True)

    tf = textbox(s, M, Inches(5.9), Inches(9.0), Inches(0.5))
    write(tf, "Kalai  ·  Xin Rong  ·  Su Myat",
          font=MONO, size=12, color=PANEL_FAINT, spacing=0.1, first=True)
    logos(s, Inches(6.45), dark=True)
    icon_building(s, W - M - Inches(2.9), Inches(2.0), Inches(2.6), PANEL_2)

    # ============================================================ 2 — the risk
    s = slide_base(prs)
    eyebrow(s, "The problem")
    tf = textbox(s, M, Inches(0.95), Inches(10.6), Inches(1.1))
    write(tf, "Construction project abandonment is a regional risk.",
          font=DISPLAY, size=34, color=INK_0, line=1.1, first=True)
    tf = textbox(s, M, Inches(1.75), Inches(9.6), Inches(0.5))
    write(tf, "49 of 581 major ASEAN infrastructure projects were cancelled or distressed "
              "after reaching financial closure.",
          font=SANS, size=14.5, color=INK_3, line=1.45, first=True)

    stats = [("8.4%", "cancelled or distressed", RESTRICTED),
             ("581", "major projects tracked", INK_0),
             ("49", "did not get built", RESTRICTED)]
    colw = (CW - Inches(0.8)) / 3
    y = Inches(2.75)
    for i, (big, label, colour) in enumerate(stats):
        x = M + i * (colw + Inches(0.4))
        box(s, x, y, colw, Inches(1.5), fill=PAPER_0, line_color=PAPER_3)
        box(s, x, y, colw, Inches(0.05), fill=colour)
        tf = textbox(s, x + Inches(0.3), y + Inches(0.32), colw - Inches(0.6), Inches(0.7))
        write(tf, big, font=DISPLAY, size=44, color=colour, line=1.0, first=True)
        tf2 = textbox(s, x + Inches(0.3), y + Inches(1.05), colw - Inches(0.6), Inches(0.35))
        write(tf2, label, font=SANS, size=12.5, color=INK_3, first=True)

    box(s, M, Inches(4.6), CW, Inches(1.35), fill=PANEL_0)
    tf = textbox(s, M + Inches(0.35), Inches(4.85), CW - Inches(0.7), Inches(1.0))
    write(tf, "When progress and the use of funds cannot be verified, the client carries the risk.",
          font=DISPLAY, size=19, color=PANEL_INK, line=1.35, first=True, space_after=6)
    write(tf, "Kirim links every payment to verified progress, and leaves a record of how the "
              "builder actually performed.",
          font=SANS, size=13.5, color=PANEL_MUTED, line=1.4)

    tf = textbox(s, M, Inches(6.15), CW, Inches(0.3))
    write(tf, "Source: OECD, Southeast Asia Investment Policy Perspectives",
          font=MONO, size=9, color=INK_4, first=True)
    footer(s, "KIRIM", "01 / PROBLEM")

    # ============================================================ 3 — it is happening now
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "Not a hypothetical", color=PANEL_FAINT)

    box(s, M, Inches(1.5), Inches(0.05), Inches(2.6), fill=CONTESTED)
    tf = textbox(s, M + Inches(0.45), Inches(1.5), Inches(10.2), Inches(2.2))
    write(tf, "Uncompleted after years, Johor Bahru's ‘sick’ and "
              "‘abandoned’ property projects leave buyers in limbo",
          font=DISPLAY, size=33, color=PANEL_INK, line=1.18, first=True, space_after=14)
    write(tf, "They bought property in Johor Bahru hoping to live or retire comfortably in "
              "Malaysia. After paying hefty sums, some are still waiting for their dream homes "
              "to be completed.",
          font=SANS, size=15, color=PANEL_MUTED, line=1.5)

    tf = textbox(s, M + Inches(0.45), Inches(4.5), Inches(10.2), Inches(0.4))
    write(tf, "Channel NewsAsia", font=MONO, size=10, color=CONTESTED, spacing=0.14, first=True)

    rule(s, M, Inches(5.3), CW, color=PANEL_2)
    tf = textbox(s, M, Inches(5.6), Inches(10.4), Inches(0.8))
    write(tf, "Buyers pay first and find out last. Every one of these started as a payment "
              "somebody released against work that was never finished.",
          font=DISPLAY, size=17, color=CONTESTED, italic=True, line=1.4, first=True)
    footer(s, "KIRIM", "02 / EVIDENCE", color=PANEL_FAINT, line=PANEL_2)

    # ============================================================ 4 — the question
    s = slide_base(prs, CONTESTED)
    box(s, 0, Inches(1.9), W, Inches(3.4), fill=PANEL_0)
    tf = textbox(s, M, Inches(2.5), CW, Inches(2.2), align=PP_ALIGN.CENTER)
    write(tf, "How might we make construction payments in ASEAN safer",
          font=DISPLAY, size=32, color=PANEL_INK, line=1.3, first=True, align=PP_ALIGN.CENTER)
    write(tf, "by releasing funds only when verified progress is achieved?",
          font=DISPLAY, size=32, color=SAGE, line=1.3, align=PP_ALIGN.CENTER)
    icon_shield(s, (W - Inches(1.2)) / 2, Inches(5.7), Inches(1.2), PANEL_0)

    # ============================================================ 5 — how it works today
    s = slide_base(prs)
    eyebrow(s, "How it works today")
    tf = textbox(s, M, Inches(0.95), Inches(10.6), Inches(0.8))
    write(tf, "Money moves first. Verification comes later, if at all.",
          font=DISPLAY, size=31, color=INK_0, line=1.1, first=True)

    steps = [
        ("Developer", "Sells the units before they are completed", "neutral"),
        ("Buyer", "Starts paying before the building is finished", "neutral"),
        ("Bank", "Mortgage money is paid out well before completion", "neutral"),
        ("Escrow account", "Holds the funds, but the controls are fragmented", "neutral"),
        ("Developer", "Has wide scope to use the cash as general funding", "bad"),
        ("Buyer", "Carries the risk if construction fails", "bad"),
        ("Result", "Work stalls after the buyer has already paid", "bad"),
    ]
    fw, fh, gap = Inches(6.4), Inches(0.62), Inches(0.135)
    fx = M + Inches(0.4)
    y = Inches(1.85)
    for who, what, tone in steps:
        flow_step(s, fx, y, fw, fh, who, what, tone=tone)
        y += fh + gap

    icon_clock(s, fx + fw + Inches(0.9), Inches(2.4), Inches(1.3), PAPER_3)
    tf = textbox(s, fx + fw + Inches(0.75), Inches(4.0), Inches(3.4), Inches(2.2))
    write(tf, "Every safeguard here depends on somebody checking, and nobody is paid to check "
              "often enough.",
          font=SANS, size=13.5, color=INK_3, line=1.5, first=True, space_after=10)
    write(tf, "By the time it is obvious, the money has gone.",
          font=DISPLAY, size=16, color=QUARANTINED, italic=True, line=1.35)
    footer(s, "KIRIM · CURRENT STATE", "03 / TODAY")

    # ============================================================ 6 — with Kirim
    s = slide_base(prs)
    eyebrow(s, "The new flow")
    tf = textbox(s, M, Inches(0.95), Inches(10.6), Inches(0.8))
    write(tf, "The same sale. The money just cannot leave until the work exists.",
          font=DISPLAY, size=29, color=INK_0, line=1.1, first=True)

    steps = [
        ("Developer", "Sells the units before they are completed", "same", "unchanged"),
        ("Buyer", "Pays a deposit up front", "same", "unchanged"),
        ("Bank", "Puts the full loan into escrow, not into the developer's account", "good", None),
        ("The escrow", "Locked on a public ledger. Nobody can move it, including us", "good", None),
        ("The agent", "Buys independent checks and reconciles them against the agreed drawings", "good", None),
        ("Release", "Conforming evidence pays in about four seconds. Anything else holds", "good", None),
        ("Refund", "A stage that is never built returns the money on its own. No lawsuit", "good", None),
    ]
    y = Inches(1.85)
    for who, what, tone, tag in steps:
        flow_step(s, fx, y, fw, fh, who, what, tone=tone, tag=tag)
        y += fh + gap

    icon_lock(s, fx + fw + Inches(0.9), Inches(2.4), Inches(1.3), TRUSTED_TINT)
    tf = textbox(s, fx + fw + Inches(0.75), Inches(4.0), Inches(3.4), Inches(2.2))
    write(tf, "Nothing about how the market sells or how buyers pay has to change. Only the "
              "destination of the money on day one.",
          font=SANS, size=13.5, color=INK_3, line=1.5, first=True, space_after=10)
    write(tf, "Two steps unchanged. Five that matter.",
          font=DISPLAY, size=16, color=TRUSTED, italic=True, line=1.35)
    footer(s, "KIRIM · WITH KIRIM", "04 / THE CHANGE")

    # ============================================================ 7 — benefits
    s = slide_base(prs)
    eyebrow(s, "What changes for everyone")
    tf = textbox(s, M, Inches(0.95), Inches(10.4), Inches(0.8))
    write(tf, "Every number here came off the ledger, not a projection.",
          font=DISPLAY, size=31, color=INK_0, line=1.1, first=True)

    rows = [
        ("Time to release a stage", "weeks of chasing", "about 4 seconds"),
        ("Cost of holding the money", "3–5% to an escrow agent", "0.8%"),
        ("Cost of checking the work", "a site visit and a surveyor", "US$0.48"),
        ("Builder is paid", "on 30 to 60 day terms", "on presentation"),
        ("If a stage is never built", "dispute, or write it off", "refunded automatically"),
        ("If a claim is rejected", "the job stops, lawyers start", "fix it and resubmit, same escrow"),
        ("Builder's reputation", "a folder of photos and hearsay", "on their own ledger account"),
    ]
    y = Inches(2.05)
    rule(s, M, y - Inches(0.14), CW, color=PAPER_3)
    for label, before, after in rows:
        tf = textbox(s, M, y, Inches(3.5), Inches(0.4))
        write(tf, label, font=SANS, size=13.5, color=INK_0, bold=True, first=True)
        tf = textbox(s, M + Inches(3.7), y + Inches(0.02), Inches(3.6), Inches(0.4))
        write(tf, before, font=SANS, size=13, color=INK_4, first=True)
        tf = textbox(s, M + Inches(7.5), y + Inches(0.02), Inches(4.0), Inches(0.4),
                     align=PP_ALIGN.RIGHT)
        write(tf, after, font=SANS, size=13.5, color=TRUSTED, bold=True, first=True,
              align=PP_ALIGN.RIGHT)
        y += Inches(0.5)
        rule(s, M, y - Inches(0.12), CW, color=PAPER_2)

    box(s, M, Inches(5.75), CW, Inches(0.95), fill=TRUSTED_TINT)
    tf = textbox(s, M + Inches(0.3), Inches(5.95), CW - Inches(0.6), Inches(0.7))
    write(tf, "Three of the six stages in our demo end with no payment at all.",
          font=DISPLAY, size=17, color=TRUSTED, italic=True, line=1.3, first=True, space_after=3)
    write(tf, "A payment that correctly does not happen is what makes an autonomous "
              "payment system worth trusting.", font=SANS, size=12.5, color=INK_2, line=1.35)
    footer(s, "KIRIM", "05 / BENEFITS")

    # ============================================================ 8 — the demo
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "The demo", color=PANEL_FAINT)
    tf = textbox(s, M, Inches(0.95), Inches(10.4), Inches(0.9))
    write(tf, "Watch a stage pay itself. Then watch one refuse to.",
          font=DISPLAY, size=32, color=PANEL_INK, line=1.1, first=True)

    beats = [
        ("1", "Foundations", "Evidence checks out. Nobody approves it. Paid in seconds.", SAGE),
        ("2", "Plumbing and interiors", "Recycled photo, edited photo, over-billed by US$20,000, "
         "72% of the agreed scope built. Refused, and the owner confirms it.", QUARANTINED),
        ("3", "The same stage again", "Defect repaired, bill corrected, real photographs. "
         "The same locked money is released. Nobody pays twice.", SAGE),
    ]
    y = Inches(2.15)
    for n, what, outcome, colour in beats:
        box(s, M, y, CW, Inches(1.05), fill=PANEL_1, line_color=PANEL_2)
        box(s, M, y, Inches(0.05), Inches(1.05), fill=colour)
        tf = textbox(s, M + Inches(0.35), y + Inches(0.2), Inches(0.5), Inches(0.6))
        write(tf, n, font=DISPLAY, size=26, color=colour, line=1.0, first=True)
        tf = textbox(s, M + Inches(1.0), y + Inches(0.19), Inches(10.0), Inches(0.75))
        write(tf, what, font=SANS, size=14.5, color=PANEL_INK, bold=True, first=True, space_after=3)
        write(tf, outcome, font=SANS, size=12.5, color=PANEL_MUTED, line=1.35)
        y += Inches(1.25)

    rule(s, M, Inches(6.05), CW, color=PANEL_2)
    tf = textbox(s, M, Inches(6.25), Inches(8.6), Inches(0.5))
    write(tf, "Real transactions on the XRP Ledger testnet. Every hash is clickable.",
          font=MONO, size=11, color=PANEL_FAINT, first=True)
    logos(s, Inches(6.2), dark=True)
    footer(s, "KIRIM", "06 / DEMO", color=PANEL_FAINT, line=PANEL_2)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KIRIM-pitch.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    with open(os.path.join(here, "hashes.json"), encoding="utf-8") as f:
        hashes = [(h["label"], h["hash"]) for h in json.load(f)]
    print("wrote", build(hashes))
