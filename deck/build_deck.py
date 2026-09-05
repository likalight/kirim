"""
KIRIM pitch deck — built in the Retia Labs design system.

Palette, structure and voice are Retia's: cream and ink, hairline rules, sharp
corners, state colours, monospaced identifiers, no gradients and no soft cards.
Typefaces are the closest system-safe analogues of Newsreader / Archivo /
IBM Plex Mono, so the deck renders correctly on a judge's machine rather than
silently substituting.

    python deck/build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import json
import os

# ---------------------------------------------------------------- Retia tokens
PAPER_0 = RGBColor(0xFA, 0xF6, 0xEE)
PAPER_1 = RGBColor(0xF1, 0xEC, 0xE1)
PAPER_2 = RGBColor(0xE8, 0xE2, 0xD4)
PAPER_3 = RGBColor(0xDE, 0xD7, 0xC6)
INK_0 = RGBColor(0x17, 0x13, 0x0D)
INK_2 = RGBColor(0x51, 0x4A, 0x3C)
INK_3 = RGBColor(0x6A, 0x63, 0x56)
INK_4 = RGBColor(0x8C, 0x84, 0x67)
PANEL_0 = RGBColor(0x0F, 0x0D, 0x09)
PANEL_1 = RGBColor(0x1A, 0x17, 0x12)
PANEL_2 = RGBColor(0x26, 0x21, 0x19)
PANEL_INK = RGBColor(0xD9, 0xD3, 0xC5)
PANEL_MUTED = RGBColor(0x9B, 0x94, 0x83)
PANEL_FAINT = RGBColor(0x7D, 0x76, 0x6A)

TRUSTED = RGBColor(0x2C, 0x5A, 0x4E)
CANDIDATE = RGBColor(0x4E, 0x7A, 0x5B)
STALE = RGBColor(0x8C, 0x84, 0x67)
CONTESTED = RGBColor(0xB0, 0x7A, 0x22)
QUARANTINED = RGBColor(0xA2, 0x3B, 0x2C)
RESTRICTED = RGBColor(0x7A, 0x2E, 0x2E)
SAGE = RGBColor(0x5E, 0x9B, 0x79)

TRUSTED_TINT = RGBColor(0xDC, 0xE6, 0xE0)
CONTESTED_TINT = RGBColor(0xEF, 0xE4, 0xCB)
QUARANTINED_TINT = RGBColor(0xEF, 0xDA, 0xD2)

# System-safe stand-ins for Newsreader / Archivo / IBM Plex Mono
DISPLAY = "Georgia"
SANS = "Segoe UI"
MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.85)          # page margin
CONTENT_W = W - 2 * M


# ---------------------------------------------------------------- primitives
def letterspace(run, ems):
    """python-pptx exposes no tracking; Retia labels depend on it."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(ems * 100)))


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def write(tf, text, *, font=SANS, size=15, color=INK_0, bold=False,
          italic=False, spacing=0.0, line=1.35, space_after=0, first=False,
          align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.line_spacing = line
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if spacing:
        letterspace(r, spacing)
    return p


def rule(slide, x, y, w, color=PAPER_3, weight=Pt(0.75)):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    ln.height = Emu(int(weight.emu / 2))
    return ln


def block(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    """Sharp-cornered rectangle. Retia has no rounded cards and no shadows."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = line_w
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def arrow_down(slide, x, y, h, color=INK_4):
    """A hairline connector with a small chevron. No clip-art arrows."""
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(9525), h)
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    tip = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x - Inches(0.05),
                                 y + h - Inches(0.07), Inches(0.1), Inches(0.08))
    tip.rotation = 180
    tip.fill.solid()
    tip.fill.fore_color.rgb = color
    tip.line.fill.background()
    tip.shadow.inherit = False


def slide_base(prs, bg=PAPER_1):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgshape = block(s, 0, 0, W, H, fill=bg)
    bgshape._element.getparent().remove(bgshape._element)
    s.shapes._spTree.insert(2, bgshape._element)
    return s


def eyebrow(slide, text, y=Inches(0.62), color=INK_4, x=M):
    tf = textbox(slide, x, y, CONTENT_W, Inches(0.25))
    write(tf, text.upper(), font=MONO, size=10.5, color=color, spacing=0.18, first=True)
    return tf


def heading(slide, text, y=Inches(0.95), size=38, color=INK_0, w=None, x=M):
    tf = textbox(slide, x, y, w or CONTENT_W, Inches(1.3))
    write(tf, text, font=DISPLAY, size=size, color=color, line=1.06, first=True)
    return tf


def footer(slide, left, right, color=INK_4):
    rule(slide, M, H - Inches(0.72), CONTENT_W, color=PAPER_3)
    tf = textbox(slide, M, H - Inches(0.58), CONTENT_W / 2, Inches(0.3))
    write(tf, left, font=MONO, size=9.5, color=color, spacing=0.12, first=True)
    tf2 = textbox(slide, M + CONTENT_W / 2, H - Inches(0.58), CONTENT_W / 2,
                  Inches(0.3), align=PP_ALIGN.RIGHT)
    write(tf2, right, font=MONO, size=9.5, color=color, spacing=0.12, first=True,
          align=PP_ALIGN.RIGHT)


LOGO_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "..", ".github", "logos")


def logos(slide, top, right=None):
    """SingHacks and Ripple, right-aligned, as single-colour lockups."""
    right = right or (W - M)
    rip = os.path.join(LOGO_DIR, "ripple.png")
    sing = os.path.join(LOGO_DIR, "singhacks.png")

    h_rip = Inches(0.3)
    pic = slide.shapes.add_picture(rip, Emu(0), top, height=h_rip)
    pic.left = Emu(int(right - pic.width))
    rip_left = pic.left

    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(int(rip_left - Inches(0.28))),
                                 top - Inches(0.02), Emu(9525), Inches(0.34))
    div.fill.solid(); div.fill.fore_color.rgb = PANEL_2
    div.line.fill.background(); div.shadow.inherit = False

    h_sing = Inches(0.34)
    pic2 = slide.shapes.add_picture(sing, Emu(0), top - Inches(0.02), height=h_sing)
    pic2.left = Emu(int(rip_left - Inches(0.56) - pic2.width))

    tf = textbox(slide, Emu(int(pic2.left - Inches(1.35))), top + Inches(0.03),
                 Inches(1.15), Inches(0.3), align=PP_ALIGN.RIGHT)
    write(tf, "BUILT FOR", font=MONO, size=9, color=PANEL_FAINT, spacing=0.16,
          first=True, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- the deck
def build(hashes):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ---------------------------------------------------------------- 1 title
    s = slide_base(prs, PANEL_0)
    tf = textbox(s, M, Inches(2.25), Inches(9.2), Inches(2.6))
    write(tf, "KIRIM", font=DISPLAY, size=88, color=PANEL_INK, line=1.0, first=True)
    write(tf, "Trust before you build.", font=DISPLAY, size=30, color=SAGE,
          italic=True, line=1.2, space_after=16)
    write(tf, "Milestone payments for construction, released by an agent that "
              "examines the evidence, settled on the XRP Ledger.",
          font=SANS, size=16, color=PANEL_MUTED, line=1.5)

    tf = textbox(s, M, Inches(1.5), CONTENT_W, Inches(0.3))
    write(tf, "SINGHACKS 2026   ·   RIPPLE   ·   AI-NATIVE BUSINESS ON XRPL",
          font=MONO, size=10.5, color=PANEL_FAINT, spacing=0.18, first=True)
    rule(s, M, Inches(2.0), CONTENT_W, color=PANEL_2, weight=Pt(1.0))

    tf = textbox(s, M, H - Inches(1.15), Inches(6.0), Inches(0.5))
    write(tf, "Less blind trust. More visible proof.", font=MONO, size=12,
          color=PANEL_FAINT, spacing=0.08, first=True)

    # the marks of the challenge this was built for, as single-colour lockups
    logos(s, H - Inches(1.28))

    # ---------------------------------------------------------------- 2 problem
    s = slide_base(prs)
    eyebrow(s, "The problem")
    heading(s, "Renovation runs on a deposit and a hope.")

    tf = textbox(s, M, Inches(2.05), Inches(6.4), Inches(2.4))
    write(tf, "A homeowner pays a large deposit before meaningful work exists. "
              "From that moment the contractor holds the money and the client "
              "holds the risk.", font=SANS, size=16, color=INK_2, line=1.55,
          space_after=14, first=True)
    write(tf, "The reliable contractor has the mirror problem. They finish the "
              "work and then wait — for a client who is slow, unhappy, or gone.",
          font=SANS, size=16, color=INK_2, line=1.55)

    stats = [
        ("50%", "typical deposit before\nany work is verified"),
        ("30–60d", "wait for payment on\ncompleted work"),
        ("$0", "recoverable when a\ncontractor walks"),
    ]
    x = M + Inches(7.0)
    for i, (big, label) in enumerate(stats):
        y = Inches(2.05) + i * Inches(1.28)
        tfa = textbox(s, x, y, Inches(2.0), Inches(0.7))
        write(tfa, big, font=DISPLAY, size=34, color=RESTRICTED, line=1.0, first=True)
        tfb = textbox(s, x + Inches(2.1), y + Inches(0.06), Inches(2.6), Inches(0.9))
        for j, ln in enumerate(label.split("\n")):
            write(tfb, ln, font=SANS, size=12.5, color=INK_3, line=1.35, first=(j == 0))
        rule(s, x, y + Inches(1.05), Inches(4.5), color=PAPER_3)

    footer(s, "KIRIM", "01 / PROBLEM")

    # ---------------------------------------------------------------- 3 current flow
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "How it works today", color=PANEL_FAINT)
    heading(s, "The money moves first. Everything else is hope.",
            color=PANEL_INK, size=32)

    steps = [
        ("CLIENT FINDS CONTRACTOR", "A name, a quote, a recommendation.", PANEL_MUTED),
        ("“CAN I TRUST THEM?”", "No verifiable record of past performance exists.", CONTESTED),
        ("CLIENT PAYS THE DEPOSIT", "Often half the contract value, up front.", QUARANTINED),
        ("CONTRACTOR CONTROLS THE MONEY", "Incentive to start the next job, not finish this one.", RESTRICTED),
    ]
    y = Inches(2.3)
    for i, (title, sub, accent) in enumerate(steps):
        bx = block(s, M, y, Inches(6.1), Inches(0.86), fill=PANEL_1, line_color=PANEL_2)
        bar = block(s, M, y, Inches(0.05), Inches(0.86), fill=accent)
        tf = textbox(s, M + Inches(0.32), y + Inches(0.14), Inches(5.6), Inches(0.6))
        write(tf, title, font=SANS, size=13.5, color=PANEL_INK, bold=True,
              spacing=0.04, first=True)
        write(tf, sub, font=SANS, size=11.5, color=PANEL_MUTED, line=1.3)
        if i < len(steps) - 1:
            arrow_down(s, M + Inches(0.55), y + Inches(0.86), Inches(0.26), color=PANEL_2)
        y += Inches(1.12)

    consequences = ["May prioritise other projects",
                    "May cause delays",
                    "Project may remain incomplete"]
    cx = M + Inches(6.6)
    cy = Inches(2.42)
    tf = textbox(s, cx, cy - Inches(0.34), Inches(4.6), Inches(0.3))
    write(tf, "WHICH MEANS", font=MONO, size=9.5, color=PANEL_FAINT,
          spacing=0.18, first=True)
    for i, c in enumerate(consequences):
        block(s, cx, cy + i * Inches(0.62), Inches(4.6), Inches(0.5),
              fill=None, line_color=PANEL_2)
        tf = textbox(s, cx + Inches(0.22), cy + i * Inches(0.62) + Inches(0.14),
                     Inches(4.2), Inches(0.3))
        write(tf, c, font=SANS, size=11.5, color=PANEL_MUTED, first=True)

    oy = cy + Inches(2.25)
    block(s, cx, oy, Inches(4.6), Inches(1.55), fill=None, line_color=RESTRICTED)
    tf = textbox(s, cx + Inches(0.22), oy + Inches(0.2), Inches(4.16), Inches(1.2))
    write(tf, "AND THE CLIENT IS LEFT WITH", font=MONO, size=9.5,
          color=RESTRICTED, spacing=0.16, first=True, space_after=8)
    write(tf, "crushing debt, and no roof over their head.", font=DISPLAY,
          size=19, color=PANEL_INK, italic=True, line=1.3)

    footer(s, "KIRIM  ·  CURRENT FLOW", "02 / PROBLEM", color=PANEL_FAINT)

    # ---------------------------------------------------------------- 4 solution
    s = slide_base(prs)
    eyebrow(s, "The solution")
    heading(s, "Pay on evidence, not on a promise.", size=40)

    tf = textbox(s, M, Inches(2.1), Inches(7.4), Inches(1.6))
    write(tf, "The project is broken into milestones. The client's money for "
              "each one is locked on the XRP Ledger before work starts — "
              "visible to the contractor, spendable by nobody.",
          font=SANS, size=17, color=INK_2, line=1.55, first=True)

    tf = textbox(s, M, Inches(3.5), Inches(7.4), Inches(1.6))
    write(tf, "When the contractor submits evidence, an agent buys what it "
              "needs to check it, reconciles it against the agreed scope, and "
              "releases the money in about four seconds — or does not, and says "
              "exactly why.", font=SANS, size=17, color=INK_2, line=1.55, first=True)

    bx = block(s, M + Inches(8.1), Inches(2.1), Inches(3.5), Inches(3.0),
               fill=PAPER_0, line_color=PAPER_3)
    tf = textbox(s, M + Inches(8.35), Inches(2.35), Inches(3.0), Inches(2.6))
    write(tf, "WHAT CHANGES", font=MONO, size=9.5, color=INK_4, spacing=0.18,
          first=True, space_after=12)
    for a, b in [("Deposit", "Escrowed milestone"),
                 ("Trust the name", "Read the record"),
                 ("Invoice, then wait", "Evidence, then paid"),
                 ("Dispute", "Automatic return")]:
        write(tf, a, font=SANS, size=12, color=INK_4, line=1.2)
        write(tf, b, font=SANS, size=14, color=TRUSTED, bold=True, line=1.2,
              space_after=10)

    footer(s, "KIRIM", "03 / SOLUTION")

    # ---------------------------------------------------------------- 5 new flow
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "The Kirim loop", color=PANEL_FAINT)
    heading(s, "Every completed milestone makes the next one easier to win.",
            color=PANEL_INK, size=30)

    loop = [
        ("MILESTONE AGREED", "Scope, date, amount, and what counts as proof."),
        ("CLIENT FUNDS ESCROW", "Locked on XRPL under a crypto-condition."),
        ("CONTRACTOR SUBMITS", "Photographs, delivery notes, permit reference."),
        ("AGENT BUYS ITS CHECKS", "Forensics, materials registry, site inspection — over x402."),
        ("EVIDENCE EXAMINED", "Reconciled against the agreed scope, line by line."),
        ("RELEASED IN ~4 SECONDS", "Or held, with the discrepancy named."),
        ("RECORD WRITTEN ON-LEDGER", "A credential on the contractor's own account."),
    ]
    colw = Inches(3.55)
    for i, (title, sub) in enumerate(loop):
        col = i // 4
        row = i % 4
        x = M + col * (colw + Inches(0.5))
        y = Inches(2.35) + row * Inches(1.02)
        accent = SAGE if i in (5, 6) else PANEL_2
        block(s, x, y, colw, Inches(0.82), fill=PANEL_1, line_color=accent)
        num = textbox(s, x + Inches(0.18), y + Inches(0.16), Inches(0.4), Inches(0.3))
        write(num, f"{i+1:02d}", font=MONO, size=10, color=SAGE, first=True)
        tf = textbox(s, x + Inches(0.62), y + Inches(0.12), colw - Inches(0.8), Inches(0.65))
        write(tf, title, font=SANS, size=12, color=PANEL_INK, bold=True,
              spacing=0.04, first=True)
        write(tf, sub, font=SANS, size=10.5, color=PANEL_MUTED, line=1.25)

    x3 = M + 2 * (colw + Inches(0.5))
    block(s, x3, Inches(2.35), colw, Inches(4.05), fill=None, line_color=SAGE)
    tf = textbox(s, x3 + Inches(0.28), Inches(2.62), colw - Inches(0.56), Inches(3.6))
    write(tf, "THE LOOP", font=MONO, size=9.5, color=SAGE, spacing=0.18,
          first=True, space_after=12)
    write(tf, "Verified work becomes a verified record. A verified record wins "
              "the next client. Which is the only reason a good contractor "
              "should out-earn a bad one.",
          font=DISPLAY, size=16, color=PANEL_INK, italic=True, line=1.38,
          space_after=14)
    write(tf, "The record lives on the contractor's XRPL account, not in our "
              "database. It is theirs, and it outlives us.",
          font=SANS, size=11.5, color=PANEL_MUTED, line=1.4)

    footer(s, "KIRIM  ·  NEW FLOW", "04 / SOLUTION", color=PANEL_FAINT)

    # ---------------------------------------------------------------- 6 instrument
    s = slide_base(prs)
    eyebrow(s, "How the money is held")
    heading(s, "An escrow that opens for a proof, not for a clock.", size=34)

    tf = textbox(s, M, Inches(2.0), Inches(6.2), Inches(2.4))
    write(tf, "Funds are locked with a PREIMAGE-SHA-256 crypto-condition. Only "
              "the holder of the fulfillment can release them, and Kirim holds "
              "it until the evidence conforms.",
          font=SANS, size=16, color=INK_2, line=1.55, space_after=14, first=True)
    write(tf, "If nothing is ever presented, CancelAfter returns the money to "
              "the client automatically. No dispute, no lawyer, no chasing.",
          font=SANS, size=16, color=INK_2, line=1.55)

    states = [("FUNDED", "money locked, visible to both sides", TRUSTED, TRUSTED_TINT),
              ("HELD", "evidence incomplete or contradicted", CONTESTED, CONTESTED_TINT),
              ("RELEASED", "conforming — paid in ~4 seconds", TRUSTED, TRUSTED_TINT),
              ("RETURNED", "nothing presented — clawed back", QUARANTINED, QUARANTINED_TINT)]
    x = M + Inches(6.9)
    for i, (name, sub, fg, tint) in enumerate(states):
        y = Inches(1.95) + i * Inches(0.92)
        block(s, x, y, Inches(4.7), Inches(0.74), fill=tint, line_color=None)
        tf = textbox(s, x + Inches(0.24), y + Inches(0.13), Inches(4.2), Inches(0.5))
        write(tf, name, font=MONO, size=11, color=fg, spacing=0.12, bold=True, first=True)
        write(tf, sub, font=SANS, size=11.5, color=INK_2, line=1.25)

    footer(s, "KIRIM", "05 / MECHANISM")

    # ---------------------------------------------------------------- 7 evidence
    s = slide_base(prs)
    eyebrow(s, "What the agent actually checks")
    heading(s, "A photograph cannot be examined. A photograph with a GPS fix can.",
            size=30)

    tf = textbox(s, M, Inches(2.0), Inches(11.6), Inches(0.6))
    write(tf, "Kirim does not claim to verify construction. It reconciles "
              "submitted evidence against the agreed scope and says plainly "
              "where the two disagree. The rules are deterministic; the model "
              "writes the advice, it never overturns the finding.",
          font=SANS, size=14.5, color=INK_2, line=1.5, first=True)

    rows = [("PHOTO-GEO", "Photograph taken outside the site boundary", "blocking"),
            ("PHOTO-TIME", "Timestamp precedes the milestone start", "blocking"),
            ("PHOTO-REUSED", "Byte-identical to an earlier submission", "blocking"),
            ("MATERIALS-SHORT", "Delivered quantity below the bill of quantities", "blocking"),
            ("DELIVERY-UNVERIFIED", "Delivery note absent from the supplier's records", "blocking"),
            ("INSPECT-INCOMPLETE", "Independent inspection below the release threshold", "blocking"),
            ("DEFECT-CRITICAL", "Critical defect open at inspection", "blocking"),
            ("SEQ-INCOMPLETE", "Prior milestone not yet released", "blocking"),
            ("LATE", "Submitted after the agreed date — recorded, not blocking", "advisory")]

    y = Inches(2.95)
    rule(s, M, y - Inches(0.12), CONTENT_W, color=PAPER_3)
    for code, desc, sev in rows:
        tf = textbox(s, M, y, Inches(3.1), Inches(0.3))
        write(tf, code, font=MONO, size=11, color=INK_0, first=True)
        tf = textbox(s, M + Inches(3.2), y, Inches(7.0), Inches(0.3))
        write(tf, desc, font=SANS, size=12.5, color=INK_2, first=True)
        tf = textbox(s, M + Inches(10.4), y, Inches(1.6), Inches(0.3), align=PP_ALIGN.RIGHT)
        write(tf, sev, font=MONO, size=10,
              color=RESTRICTED if sev == "blocking" else CONTESTED,
              first=True, align=PP_ALIGN.RIGHT)
        y += Inches(0.4)
        rule(s, M, y - Inches(0.09), CONTENT_W, color=RGBColor(0xE8, 0xE2, 0xD4))

    footer(s, "KIRIM", "06 / MECHANISM")

    # ---------------------------------------------------------------- 8 outcomes
    s = slide_base(prs)
    eyebrow(s, "Three outcomes, not two")
    heading(s, "“You did not send enough” and “this does not add up” "
               "are different messages.", size=30)

    cards = [("READY", "Evidence is consistent with the agreed scope. Released "
                       "autonomously below the client's ceiling.", TRUSTED, TRUSTED_TINT),
             ("MORE INFORMATION", "Nothing contradicts the scope — the submission "
                                  "is simply incomplete. Funds held. No mark against "
                                  "the record.", CONTESTED, CONTESTED_TINT),
             ("FLAGGED", "Evidence contradicts the scope. Funds held, discrepancy "
                         "named, client reviews.", QUARANTINED, QUARANTINED_TINT)]
    cw = Inches(3.65)
    for i, (name, body, fg, tint) in enumerate(cards):
        x = M + i * (cw + Inches(0.35))
        block(s, x, Inches(2.55), cw, Inches(2.6), fill=tint, line_color=None)
        block(s, x, Inches(2.55), cw, Inches(0.06), fill=fg)
        tf = textbox(s, x + Inches(0.3), Inches(2.9), cw - Inches(0.6), Inches(2.1))
        write(tf, name, font=MONO, size=11.5, color=fg, spacing=0.14, bold=True,
              first=True, space_after=12)
        write(tf, body, font=SANS, size=13.5, color=INK_2, line=1.45)

    tf = textbox(s, M, Inches(5.5), Inches(11.6), Inches(0.6))
    write(tf, "Only the third one should ever damage a contractor's track "
              "record. Most systems collapse these into a single rejection, and "
              "punish the honest contractor who forgot a photograph.",
          font=DISPLAY, size=16, color=INK_2, italic=True, line=1.45, first=True)

    footer(s, "KIRIM", "07 / MECHANISM")

    # ---------------------------------------------------------------- 9 controls
    s = slide_base(prs)
    eyebrow(s, "Autonomy, with a brake")
    heading(s, "The agent may request a payment. Only the ledger service may send one.",
            size=30)

    tf = textbox(s, M, Inches(2.15), Inches(6.2), Inches(2.2))
    write(tf, "The agent holds no wallet seed and never touches xrpl.js. Every "
              "movement of money is a request to a separate process that owns "
              "the keys, enforces the ceilings, and can refuse — and a refusal "
              "is a logged decision with a reason, not a silent no-op.",
          font=SANS, size=15.5, color=INK_2, line=1.55, first=True)

    ctrl = [("Per x402 call", "US$1.00"),
            ("Per milestone, operating", "US$5.00"),
            ("Autonomous release ceiling", "US$12,000"),
            ("Above the ceiling", "client authorises")]
    x = M + Inches(6.9)
    y = Inches(2.05)
    rule(s, x, y - Inches(0.12), Inches(4.7), color=PAPER_3)
    for label, val in ctrl:
        tf = textbox(s, x, y, Inches(3.0), Inches(0.35))
        write(tf, label, font=SANS, size=13, color=INK_2, first=True)
        tf = textbox(s, x + Inches(3.0), y, Inches(1.7), Inches(0.35), align=PP_ALIGN.RIGHT)
        write(tf, val, font=MONO, size=12.5, color=INK_0, first=True, align=PP_ALIGN.RIGHT)
        y += Inches(0.52)
        rule(s, x, y - Inches(0.11), Inches(4.7), color=PAPER_2)

    tf = textbox(s, x, y + Inches(0.2), Inches(4.7), Inches(1.4))
    write(tf, "The brief names “humans approving every agent action” as an "
              "anti-pattern. A value threshold is a safeguard, not an approval "
              "queue: small milestones settle themselves, large ones ask.",
          font=SANS, size=12.5, color=INK_3, line=1.45, first=True)

    footer(s, "KIRIM", "08 / GOVERNANCE")

    # ---------------------------------------------------------------- 10 record
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "The track record", color=PANEL_FAINT)
    heading(s, "A reputation you own, not a row in our database.",
            color=PANEL_INK, size=34)

    tf = textbox(s, M, Inches(2.15), Inches(6.0), Inches(2.4))
    write(tf, "Every released milestone issues an XLS-70 Credential to the "
              "contractor's own XRPL account, which they accept. It is keyed to "
              "the project and milestone, so it cannot be double-counted, and "
              "any future client can verify it without asking Kirim anything.",
          font=SANS, size=15.5, color=PANEL_MUTED, line=1.55, space_after=16,
          first=True)
    write(tf, "We do not claim a credential makes a contractor trustworthy. It "
              "makes their history visible.",
          font=DISPLAY, size=17, color=SAGE, italic=True, line=1.4)

    x = M + Inches(6.7)
    block(s, x, Inches(2.05), Inches(4.9), Inches(3.1), fill=PANEL_1, line_color=PANEL_2)
    tf = textbox(s, x + Inches(0.32), Inches(2.32), Inches(4.3), Inches(2.6))
    write(tf, "ON-LEDGER CREDENTIAL", font=MONO, size=9.5, color=SAGE,
          spacing=0.18, first=True, space_after=14)
    for k, v in [("type", "KIRIM:PRJ-2026-014:M1"),
                 ("subject", "rhayr2jyg…XLvv3"),
                 ("issuer", "rhSezETXu…b5aMF"),
                 ("uri", "kirim:milestone/PRJ-2026-014/"),
                 ("", "M1/demolition-and-disposal?onTime=1"),
                 ("accepted", "true")]:
        p = tf.add_paragraph()
        p.line_spacing = 1.5
        r = p.add_run()
        r.text = f"{k:<10}" if k else " " * 10
        r.font.name = MONO
        r.font.size = Pt(11)
        r.font.color.rgb = PANEL_FAINT
        r2 = p.add_run()
        r2.text = v
        r2.font.name = MONO
        r2.font.size = Pt(11)
        r2.font.color.rgb = PANEL_INK

    footer(s, "KIRIM  ·  XLS-70 CREDENTIALS", "09 / TRACK RECORD", color=PANEL_FAINT)

    # ---------------------------------------------------------------- 11 built
    s = slide_base(prs)
    eyebrow(s, "Work done so far")
    heading(s, "Built and verified on XRPL testnet.", size=38)

    built = [
        ("Escrow with crypto-conditions", "EscrowCreate / EscrowFinish / EscrowCancel, hand-encoded PREIMAGE-SHA-256", "done"),
        ("x402, both sides", "402 challenge, payment, on-ledger verification before a byte is served", "done"),
        ("Evidence examination", "Deterministic rules over photos, deliveries, permits, inspection", "done"),
        ("Spend controls", "Per-call, per-milestone, per-run ceilings enforced server-side", "done"),
        ("XLS-70 track record", "Credential issued and accepted on the contractor's account", "done"),
        ("Live decision log", "Server-sent events, every decision with its reason and hash", "done"),
        ("RLUSD settlement", "Issuer wired, trustlines placed — faucet funding outstanding", "partial"),
        ("Agent credit, DEX leg", "Draw-down instead of pre-funding; FX executed on the XRPL DEX", "next"),
    ]
    y = Inches(2.15)
    rule(s, M, y - Inches(0.12), CONTENT_W, color=PAPER_3)
    for name, desc, state in built:
        colour = {"done": TRUSTED, "partial": CONTESTED, "next": INK_4}[state]
        block(s, M, y + Inches(0.08), Inches(0.06), Inches(0.3), fill=colour)
        tf = textbox(s, M + Inches(0.24), y, Inches(3.6), Inches(0.4))
        write(tf, name, font=SANS, size=13.5, color=INK_0, bold=True, first=True)
        tf = textbox(s, M + Inches(3.95), y + Inches(0.02), Inches(6.6), Inches(0.4))
        write(tf, desc, font=SANS, size=12.5, color=INK_2, first=True)
        tf = textbox(s, M + Inches(10.7), y + Inches(0.02), Inches(1.3), Inches(0.35),
                     align=PP_ALIGN.RIGHT)
        write(tf, state, font=MONO, size=10, color=colour, spacing=0.1,
              first=True, align=PP_ALIGN.RIGHT)
        y += Inches(0.52)
        rule(s, M, y - Inches(0.11), CONTENT_W, color=PAPER_2)

    footer(s, "KIRIM", "10 / BUILT")

    # ---------------------------------------------------------------- 12 proof
    s = slide_base(prs, PANEL_0)
    eyebrow(s, "Proof", color=PANEL_FAINT)
    heading(s, "Every payment on this slide is real.", color=PANEL_INK, size=36)

    tf = textbox(s, M, Inches(1.95), Inches(11.6), Inches(0.5))
    write(tf, "XRP Ledger Testnet. Supply is simulated — the contractor, the "
              "photographs, the inspector. The ledger is not.",
          font=SANS, size=14, color=PANEL_MUTED, first=True)

    y = Inches(2.75)
    rule(s, M, y - Inches(0.14), CONTENT_W, color=PANEL_2)
    for label, h in hashes:
        tf = textbox(s, M, y, Inches(3.5), Inches(0.32))
        write(tf, label, font=SANS, size=12.5, color=PANEL_INK, first=True)
        tf = textbox(s, M + Inches(3.6), y + Inches(0.02), Inches(8.0), Inches(0.32))
        write(tf, h, font=MONO, size=10.5, color=SAGE, first=True)
        y += Inches(0.46)
        rule(s, M, y - Inches(0.12), CONTENT_W, color=PANEL_2)

    tf = textbox(s, M, H - Inches(1.55), Inches(11.6), Inches(0.4))
    write(tf, "testnet.xrpl.org/transactions/<hash>", font=MONO, size=11,
          color=PANEL_FAINT, first=True)

    footer(s, "KIRIM  ·  XRPL TESTNET", "11 / PROOF", color=PANEL_FAINT)

    # ---------------------------------------------------------------- 13 model
    s = slide_base(prs)
    eyebrow(s, "The business")
    heading(s, "We charge for the release, because that is where the risk is.",
            size=32)

    tf = textbox(s, M, Inches(2.15), Inches(6.0), Inches(2.4))
    write(tf, "Kirim takes 0.8% of each released milestone, paid by the client "
              "at the moment of release. A licensed escrow agent charges 3–5% "
              "and takes days; a lawyer's stakeholder account costs more than "
              "the milestone.",
          font=SANS, size=15.5, color=INK_2, line=1.55, space_after=16, first=True)
    write(tf, "The evidence checks are bought per milestone at cents apiece, so "
              "the cost of assurance no longer scales with a site visit.",
          font=SANS, size=15.5, color=INK_2, line=1.55)

    x = M + Inches(6.9)
    rows = [("Escrow agent, today", "3–5%", "days"),
            ("Kirim", "0.8%", "~4 seconds"),
            ("Evidence checks", "US$0.48", "per milestone"),
            ("On a S$50,000 renovation", "S$400", "total fee")]
    y = Inches(2.05)
    rule(s, x, y - Inches(0.12), Inches(4.7), color=PAPER_3)
    for a, b, c in rows:
        tf = textbox(s, x, y, Inches(2.6), Inches(0.4))
        write(tf, a, font=SANS, size=13, color=INK_2, first=True)
        tf = textbox(s, x + Inches(2.6), y, Inches(1.1), Inches(0.4), align=PP_ALIGN.RIGHT)
        write(tf, b, font=MONO, size=13, color=INK_0, first=True, align=PP_ALIGN.RIGHT)
        tf = textbox(s, x + Inches(3.75), y + Inches(0.03), Inches(0.95), Inches(0.4),
                     align=PP_ALIGN.RIGHT)
        write(tf, c, font=SANS, size=11, color=INK_4, first=True, align=PP_ALIGN.RIGHT)
        y += Inches(0.58)
        rule(s, x, y - Inches(0.12), Inches(4.7), color=PAPER_2)

    footer(s, "KIRIM", "12 / BUSINESS")

    # ---------------------------------------------------------------- 14 close
    s = slide_base(prs, PANEL_0)
    tf = textbox(s, M, Inches(2.2), Inches(9.6), Inches(3.0))
    write(tf, "Remove the agent and you need a\nsite visit for every payment.",
          font=DISPLAY, size=34, color=PANEL_INK, line=1.2, first=True,
          space_after=18)
    write(tf, "Remove autonomous payment and the\nescrow is just an invoice again.",
          font=DISPLAY, size=34, color=SAGE, line=1.2)

    rule(s, M, Inches(5.6), CONTENT_W, color=PANEL_2, weight=Pt(1.0))
    tf = textbox(s, M, Inches(5.85), Inches(7.0), Inches(0.6))
    write(tf, "KIRIM  ·  LESS BLIND TRUST. MORE VISIBLE PROOF.", font=MONO,
          size=12, color=PANEL_FAINT, spacing=0.18, first=True)
    logos(s, Inches(5.78))

    out = os.path.join(os.path.dirname(__file__), "KIRIM-pitch-deck.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    here = os.path.dirname(__file__)
    with open(os.path.join(here, "hashes.json"), encoding="utf-8") as f:
        hashes = [(h["label"], h["hash"]) for h in json.load(f)]
    path = build(hashes)
    print("wrote", path)
